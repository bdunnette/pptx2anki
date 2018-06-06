#!/usr/bin/env python3

import sys, os
import genanki
from random import randrange
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

deck_name = os.path.splitext(os.path.basename(sys.argv[1]))[0]

my_model = genanki.Model(
  randrange(9999999999),
  'Simple Model',
  fields=[
    {'name': 'Question'},
    {'name': 'Answer'},
  ],
  templates=[
    {
      'name': 'Card 1',
      'qfmt': '{{Question}}',
      'afmt': '{{FrontSide}}<hr id="answer">{{Answer}}',
    },
  ])

my_deck = genanki.Deck(
  # Generate a large random deck ID to prevent decks 'overlapping' on import
  randrange(9999999999),
  deck_name)

my_package = genanki.Package(my_deck)
my_package.media_files = []

prs = Presentation(sys.argv[1])
# print(prs)
print("%s slides found..." % len(prs.slides))
slide_info = [dict() for i in range(len(prs.slides))]
# Loop through slides, adding them to slide_info dictionary, so that we can look ahead to next slide
for slide_index, slide in enumerate(prs.slides):
    note_text = ""
    # Compile all text on this slide's notes_slide
    for ns in slide.notes_slide.shapes:
        if not str(ns.text).isnumeric():
            note_text += str(ns.text)
    slide_info[slide_index]['text'] = note_text.strip()
    # Extract image files so that we can import them into our Anki package
    for shape in slide.shapes:
        print(shape.shape_type)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_file = "%s.%s" % (slide_index, shape.image.ext)
            with open(image_file, 'wb') as f:
                f.write(shape.image.blob)
            slide_info[slide_index]['image'] = image_file
    print(slide_info[slide_index])

# Loop through slides in slide_info
for idx, slide in enumerate(slide_info):
    print("Processing slide %s" % idx)
    if 'image' in slide and not 'processed' in slide:
        # print(slide)
        slide_text = slide['text'].strip()
        print("Slide text: %s" % slide_text)
        next_slide = slide_info[idx+1]
        next_slide_text = next_slide['text'].strip()
        print("Next slide text: %s" % next_slide_text)
        back_slide_exists = ((slide_text == next_slide_text) or (slide_text == next_slide_text.rstrip(".")) or (slide_text.rstrip(".") == next_slide_text))
        print("Back slide found: %s" % back_slide_exists)
        print("Slide notes: %s" % slide_text)
        # If next slide has same note text as this one, add that slide's image as the "back" of our card
        if (back_slide_exists) and ('image' in next_slide):
            back_text = "<img src='%s'><p>%s</p>" % (next_slide['image'],slide_text)
            my_package.media_files.append(next_slide['image'])
            slide_info[idx+1]['processed'] = True
        # Otherwise, just put the note text on the card back
        else:
            back_text = slide_text
        # Add note to deck
        my_note = genanki.Note(
          model=my_model,
          sort_field=idx,
          fields=["<img src='%s'>" % slide['image'],back_text])
        my_deck.add_note(my_note)
        my_package.media_files.append(slide['image'])
        slide_info[idx]['processed'] = True

# Export deck to .apkg file
my_package.write_to_file('%s.apkg' % deck_name)
